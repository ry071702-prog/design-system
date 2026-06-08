#!/usr/bin/env python3
"""
build_pptx.py — スライド構成ジェネレーターの出力を、サイバーエージェント公式
テンプレートをベースにした .pptx に書き出す (ローカル専用ツール)。

仕組み: 公式テンプレ pptx を開き、必要な作例スライド(表紙/AGENDA/本文)を
"複製"して文字を差し替え、元の52枚のガイドラインスライドを削除して保存する。
これでテーマ色(CAグリーン)・フォント・ロゴ・レイアウトを完全に継承できる。

使い方:
  # Dify を呼んで構成を作り、そのまま pptx 化
  python3 tools/build_pptx.py --topic "社内ナレッジ検索へのRAG導入" \
      --theme editorial --audience 経営層 --goal 予算承認 --constraints "全10枚以内" --open

  # 既に手元にある構成 Markdown から pptx 化 (Dify を呼ばない)
  python3 tools/build_pptx.py --from out/structure.md --open

設定: config.local.json の difyBaseUrl / difyAppKey / pptxTemplate を使う。
"""
import argparse, json, os, re, sys, copy, datetime, subprocess, urllib.request

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
CONFIG = os.path.join(ROOT, "config.local.json")
DEFAULT_TEMPLATE = os.path.expanduser(
    "~/Downloads/会社資料/ガイドライン・ブランド/"
    "サイバーエージェント公式スライドガイドライン _ テンプレート[共有用].pptx"
)
# 複製元にする作例スライドの index (テンプレ内)
SRC_COVER, SRC_AGENDA, SRC_CONTENT = 21, 22, 26

RT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
RT_LAYOUT = RT + "/slideLayout"
RT_NOTES = RT + "/notesSlide"
ACCENT_GREEN = "298737"  # theme accent1 (CA グリーン)
BODY_FONT = "M PLUS 1p"  # テンプレ推奨フォント


# ---------- 設定 ----------
def load_config():
    if not os.path.exists(CONFIG):
        sys.exit("config.local.json がありません。config.example.json をコピーして設定してください。")
    with open(CONFIG, encoding="utf-8") as f:
        return json.load(f)


# ---------- Dify ----------
def call_dify(cfg, inputs):
    base = cfg["difyBaseUrl"].rstrip("/")
    body = json.dumps({"inputs": inputs, "response_mode": "blocking",
                       "user": "build-pptx-cli"}).encode("utf-8")
    req = urllib.request.Request(
        base + "/workflows/run", data=body,
        headers={"Authorization": "Bearer " + cfg["difyAppKey"],
                 "Content-Type": "application/json"})
    with urllib.request.urlopen(req, timeout=300) as r:
        data = json.loads(r.read().decode("utf-8"))
    outputs = (data.get("data") or {}).get("outputs") or {}
    for k in ("skeleton", "text", "result", "output"):
        if isinstance(outputs.get(k), str) and outputs[k].strip():
            return outputs[k]
    vals = [v for v in outputs.values() if isinstance(v, str) and v.strip()]
    if not vals:
        sys.exit("Dify の出力が空でした: " + json.dumps(data, ensure_ascii=False)[:300])
    return vals[0]


# ---------- 構成 Markdown のパース (LLM 出力の揺れに頑健) ----------
def _clean(s):
    return s.replace("**", "").replace("＊＊", "").strip(" 　")


def parse_deck(md):
    title, slides, cur = None, [], None
    for raw in md.splitlines():
        line = raw.strip()
        if not line:
            continue
        hm = re.match(r"^(#{1,6})\s*(.+)$", line)
        if hm:
            level, text = len(hm.group(1)), _clean(hm.group(2))
            if level == 1 and title is None:
                title = re.split(r"[—\-―〜~｜|:：]", text)[0].strip() or text
                continue
            if level >= 2 or "スライド" in text:
                if cur:
                    slides.append(cur)
                role = re.sub(r"^スライド\s*\d+\s*[:：.\s　]+", "", text)
                role = re.sub(r"^\d+[\.\s　]+", "", role).strip()
                cur = {"role": role, "layout": "", "elements": [], "hint": ""}
                continue
        if cur is None:
            continue
        body = _clean(re.sub(r"^[-*・•‣▪]\s*", "", line))
        m = re.match(r"^(?:レイアウト|構図)\s*[:：]\s*(.+)$", body)
        if m:
            cur["layout"] = m.group(1).strip()
            continue
        m = re.match(r"^(?:載せる要素|要素|配置)\s*[:：]\s*(.+)$", body)
        if m:
            cur["elements"] = [x.strip() for x in re.split(r"[／/、,]\s*", m.group(1)) if x.strip()]
            continue
        m = re.match(r"^(?:記入ヒント|ヒント)\s*[:：]\s*(.+)$", body)
        if m:
            cur["hint"] = m.group(1).strip()
            continue
    if cur:
        slides.append(cur)
    return (title or "資料タイトル"), slides


# ---------- pptx 複製ユーティリティ ----------
def duplicate_slide(prs, index):
    from pptx.oxml.ns import qn  # noqa
    src = prs.slides[index]
    new = prs.slides.add_slide(src.slide_layout)
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)
    remap = {}
    for rId, rel in src.part.rels.items():
        if rel.reltype in (RT_LAYOUT, RT_NOTES):
            continue
        if rel.is_external:
            remap[rId] = new.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            remap[rId] = new.part.relate_to(rel.target_part, rel.reltype)
    rns = "{" + RT + "}"
    for shp in src.shapes:
        el = copy.deepcopy(shp._element)
        for node in el.iter():
            for a, v in list(node.attrib.items()):
                if a.startswith(rns) and v in remap:
                    node.set(a, remap[v])
        new.shapes._spTree.append(el)
    return new


def delete_slides(prs, sld_ids):
    """sld_ids: 削除する <p:sldId> 要素のリスト。rel も落として保存時に剪定させる。"""
    from pptx.oxml.ns import qn
    lst = prs.slides._sldIdLst
    for sldId in sld_ids:
        rId = sldId.get(qn("r:id"))
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass
        lst.remove(sldId)


def find_shape(slide, text):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == text:
            return sh
    return None


def set_text(shape, text):
    """先頭 run の書式を保ったままテキストを差し替え (余分な run/段落は削除)。"""
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = text
        for r in p0.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        p0.add_run().text = text
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)


# ---------- 各スライド組み立て ----------
def build_cover(prs, title, subtitle, date_str):
    s = duplicate_slide(prs, SRC_COVER)
    # 値が空のもの(subtitle 等)はテンプレのプレースホルダ文言を残して記入を促す
    for placeholder, value in (("タイトル1段の場合", title),
                               ("サブタイトル入力", subtitle),
                               ("2025.00.00", date_str)):
        if not value:
            continue
        sh = find_shape(s, placeholder)
        if sh:
            set_text(sh, value)
    return s


def build_agenda(prs, slides):
    from pptx.oxml.ns import qn
    s = duplicate_slide(prs, SRC_AGENDA)
    # 目次テキストの入った図形を探して、役割一覧で置き換える
    target = None
    for sh in s.shapes:
        if sh.has_text_frame and "目次テキスト" in sh.text_frame.text:
            target = sh
            break
    if target:
        tf = target.text_frame
        base_p = tf.paragraphs[0]
        for extra in tf.paragraphs[1:]:
            extra._p.getparent().remove(extra._p)
        roles = [sl["role"] for sl in slides][:8]
        if base_p.runs:
            base_p.runs[0].text = f"01　{roles[0]}" if roles else "01"
            for r in base_p.runs[1:]:
                r._r.getparent().remove(r._r)
        for i, role in enumerate(roles[1:], start=2):
            p = tf.add_paragraph()
            r = p.add_run()
            r.text = f"{i:02d}　{role}"
    return s


COVER_KW = ("表紙", "タイトル", "cover", "title")
AGENDA_KW = ("目次", "アジェンダ", "agenda", "contents", "コンテンツ", "もくじ")


def is_cover(sl):
    return any(k in sl["role"].lower() for k in COVER_KW)


def is_agenda(sl):
    return any(k in sl["role"].lower() for k in AGENDA_KW)


def build_content(prs, idx, sl):
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    s = duplicate_slide(prs, SRC_CONTENT)
    title_sh = find_shape(s, "タイトル")
    if title_sh:
        set_text(title_sh, f"{idx:02d}　{sl['role']}")
    # 本文ガイド (枠) を新規テキストボックスで追加
    tb = s.shapes.add_textbox(Emu(640000), Emu(1650000), Emu(7850000), Emu(2900000))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"レイアウト: {sl['layout'] or '（自由）'}"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.name = BODY_FONT
    r.font.color.rgb = RGBColor.from_string(ACCENT_GREEN)
    if sl["elements"]:
        head = tf.add_paragraph()
        rh = head.add_run(); rh.text = "載せる要素:"
        rh.font.size = Pt(12); rh.font.bold = True; rh.font.name = BODY_FONT
        for el in sl["elements"]:
            pe = tf.add_paragraph()
            re_ = pe.add_run(); re_.text = "・" + el
            re_.font.size = Pt(12); re_.font.name = BODY_FONT
    if sl["hint"]:
        ph = tf.add_paragraph()
        rp = ph.add_run(); rp.text = "記入ヒント: " + sl["hint"]
        rp.font.size = Pt(11); rp.font.name = BODY_FONT
        rp.font.color.rgb = RGBColor.from_string("888888")
    return s


# ---------- メイン ----------
def slugify(s):
    s = re.sub(r"[\\/:*?\"<>|]+", "", s).strip()
    return (s[:40] or "proposal")


def main():
    ap = argparse.ArgumentParser(description="スライド構成 → CAテンプレ pptx")
    ap.add_argument("--topic", help="提案テーマ (指定すると Dify を呼ぶ)")
    ap.add_argument("--theme", default="editorial", help="studio/editorial/focus")
    ap.add_argument("--tone", default="")
    ap.add_argument("--audience", default="")
    ap.add_argument("--goal", default="")
    ap.add_argument("--constraints", default="")
    ap.add_argument("--count", default="", help="スライド枚数 (任意)")
    ap.add_argument("--outline", default="", help="各スライドの大まかな内容/全体イメージ (任意)")
    ap.add_argument("--from", dest="from_file", help="構成 Markdown ファイルから作る (Dify を呼ばない)")
    ap.add_argument("--out", help="出力 pptx パス")
    ap.add_argument("--template", help="ベースにする pptx (既定: config.local.json か Downloads の公式テンプレ)")
    ap.add_argument("--open", action="store_true", help="生成後に開く (macOS)")
    args = ap.parse_args()

    cfg = load_config()
    template = args.template or cfg.get("pptxTemplate") or DEFAULT_TEMPLATE
    if not os.path.exists(template):
        sys.exit(f"テンプレが見つかりません: {template}\n--template でパス指定するか config.local.json の pptxTemplate を設定してください。")

    if args.from_file:
        with open(args.from_file, encoding="utf-8") as f:
            md = f.read()
    elif args.topic:
        print("Dify で構成を生成中…", file=sys.stderr)
        md = call_dify(cfg, {"mode": "deck", "theme": args.theme, "tone": args.tone,
                             "topic": args.topic, "audience": args.audience,
                             "goal": args.goal, "constraints": args.constraints, "idea": "",
                             "count": args.count, "outline": args.outline})
    else:
        sys.exit("--topic か --from のどちらかを指定してください。")

    title, slides = parse_deck(md)
    if not slides:
        sys.exit("構成を解析できませんでした。Markdown の見出し形式を確認してください。\n--- 入力 ---\n" + md[:500])
    print(f"タイトル: {title} / スライド {len(slides)} 枚", file=sys.stderr)

    from pptx import Presentation
    prs = Presentation(template)
    originals = list(prs.slides._sldIdLst)  # 後で消す元の52枚

    date_str = datetime.date.today().strftime("%Y.%m.%d")
    cover_title = args.topic or title          # 表紙の主タイトルは提案テーマを優先
    cover_sub = args.goal or ""                 # 空ならテンプレの記入プレースホルダを残す
    content_slides = [s for s in slides if not is_cover(s) and not is_agenda(s)]
    # 役割ごとにテンプレを選択。表紙が無ければ先頭に補う (枚数は構成に忠実)
    if not any(is_cover(s) for s in slides):
        build_cover(prs, cover_title, cover_sub, date_str)
    content_idx = 0
    for sl in slides:
        if is_cover(sl):
            build_cover(prs, cover_title, cover_sub, date_str)
        elif is_agenda(sl):
            build_agenda(prs, content_slides)
        else:
            content_idx += 1
            build_content(prs, content_idx, sl)

    delete_slides(prs, originals)

    final_count = len(prs.slides._sldIdLst)
    out = args.out or os.path.expanduser(f"~/Downloads/{slugify(title)}_v0.1.pptx")
    prs.save(out)
    size_mb = os.path.getsize(out) / 1024 / 1024
    print(f"[OK] 出力: {out}  ({size_mb:.1f} MB, {final_count} スライド)")
    if args.open:
        subprocess.run(["open", out], check=False)


if __name__ == "__main__":
    main()
