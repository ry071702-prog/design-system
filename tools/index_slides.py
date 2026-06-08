#!/usr/bin/env python3
"""
index_slides.py — これまで作成したスライド(pptx / key)を走査して、
一覧ページ (slides.html) 用の索引 data/slides.json とサムネイル
(data/slide-thumbs/*.png) を生成する (ローカル専用)。

サムネは macOS の qlmanage (QuickLook) で1枚目を描画。pptx の枚数と
タイトルは python-pptx で取得 (失敗時はファイル名で代替)。

使い方:
  python3 tools/index_slides.py            # 既定で ~/Downloads を走査
  python3 tools/index_slides.py ~/Desktop  # 走査ルートを指定 (複数可)
設定: config.local.json の slidesRoots(配列) があればそれを既定にする。
"""
import json, os, re, sys, hashlib, subprocess, datetime

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
CONFIG = os.path.join(ROOT, "config.local.json")
THUMB_DIR = os.path.join(ROOT, "data", "slide-thumbs")
OUT_JSON = os.path.join(ROOT, "data", "slides.json")
EXTS = (".pptx", ".key")
EXCLUDE = ("_削除候補", "/Library/", "node_modules", "/.Trash", os.path.basename(ROOT))


def load_roots(argv):
    if argv:
        return [os.path.expanduser(a) for a in argv]
    if os.path.exists(CONFIG):
        with open(CONFIG, encoding="utf-8") as f:
            r = (json.load(f) or {}).get("slidesRoots")
            if r:
                return [os.path.expanduser(x) for x in r]
    return [os.path.expanduser("~/Downloads")]


def find_files(roots):
    found = []
    for root in roots:
        for dirpath, dirs, files in os.walk(root):
            if any(x in dirpath for x in EXCLUDE):
                dirs[:] = []
                continue
            for fn in files:
                if fn.startswith("~$") or fn.startswith("."):
                    continue
                if fn.lower().endswith(EXTS):
                    found.append(os.path.join(dirpath, fn))
    return found


def pptx_meta(path):
    """(slide_count, title) を返す。pptx 以外/失敗時は (None, None)。"""
    if not path.lower().endswith(".pptx"):
        return None, None
    try:
        from pptx import Presentation
        prs = Presentation(path)
        count = len(prs.slides._sldIdLst)
        title = None
        for s in prs.slides:
            for sh in s.shapes:
                if sh.has_text_frame and sh.text_frame.text.strip():
                    title = sh.text_frame.text.strip().splitlines()[0][:60]
                    break
            if title:
                break
        return count, title
    except Exception:
        return None, None


def thumbnail(path, key):
    """qlmanage で 1枚目サムネを生成し、data/slide-thumbs/<key>.png の相対パスを返す。"""
    dst = os.path.join(THUMB_DIR, key + ".png")
    rel = os.path.relpath(dst, ROOT)
    if os.path.exists(dst) and os.path.getmtime(dst) >= os.path.getmtime(path):
        return rel
    tmp = os.path.join(THUMB_DIR, "_tmp")
    os.makedirs(tmp, exist_ok=True)
    for f in os.listdir(tmp):
        os.remove(os.path.join(tmp, f))
    subprocess.run(["qlmanage", "-t", "-s", "800", "-o", tmp, path],
                   stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=60)
    pngs = [f for f in os.listdir(tmp) if f.lower().endswith(".png")]
    if not pngs:
        return ""
    os.replace(os.path.join(tmp, pngs[0]), dst)
    return rel


def main():
    roots = load_roots(sys.argv[1:])
    os.makedirs(THUMB_DIR, exist_ok=True)
    files = find_files(roots)
    print(f"走査: {roots} → {len(files)} 件", file=sys.stderr)

    items = []
    for i, path in enumerate(sorted(files), 1):
        st = os.stat(path)
        key = hashlib.sha1(path.encode("utf-8")).hexdigest()[:16]
        count, title = pptx_meta(path)
        thumb = thumbnail(path, key)
        folder = os.path.basename(os.path.dirname(path))
        items.append({
            "id": key,
            "name": os.path.splitext(os.path.basename(path))[0],
            "title": title or os.path.splitext(os.path.basename(path))[0],
            "folder": folder,
            "path": path,
            "ext": os.path.splitext(path)[1].lstrip(".").lower(),
            "slides": count,
            "sizeMB": round(st.st_size / 1024 / 1024, 1),
            "mtime": datetime.date.fromtimestamp(st.st_mtime).isoformat(),
            "thumb": thumb,
        })
        if i % 10 == 0:
            print(f"  {i}/{len(files)}", file=sys.stderr)

    items.sort(key=lambda x: x["mtime"], reverse=True)
    with open(OUT_JSON, "w", encoding="utf-8") as f:
        json.dump({"generated": datetime.date.today().isoformat(),
                   "roots": roots, "slides": items}, f, ensure_ascii=False, indent=2)
    withthumb = sum(1 for x in items if x["thumb"])
    print(f"[OK] {OUT_JSON} に {len(items)} 件 (サムネ {withthumb} 件)")


if __name__ == "__main__":
    main()
